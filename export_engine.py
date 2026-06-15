#!/usr/bin/env python3
"""
导出引擎 v3 — PDF (浏览器渲染) + PPT (全板块可编辑)
"""
import json
from pathlib import Path
from datetime import datetime


# ============================================================
# PDF 导出 — Playwright 渲染 report.html
# ============================================================

def export_pdf(data_file, desktop):
    """
    生成 PDF — 使用 Playwright 渲染 report.html 完整页面
    这样 PDF 与浏览器显示完全一致，不会乱码
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return None, "playwright 未安装，请运行: pip install playwright && playwright install chromium"

    report_html = Path(data_file).parent / "report.html"
    if not report_html.exists():
        return None, "report.html 不存在"

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = Path(desktop) / f"AI行业周报_{timestamp}.pdf"

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1280, "height": 900})

        # 加载 report.html
        file_url = f"file:///{report_html.resolve().as_posix()}"
        page.goto(file_url, wait_until="networkidle", timeout=30000)

        # 等待 JS 渲染完成
        page.wait_for_timeout(3000)

        # 注入打印优化 CSS
        page.add_style_tag(content="""
            .toolbar, .nav-bar, .back-top, .toast, .filter-bar,
            .btn-trans, #statusBar { display: none !important; }
            .section { break-inside: avoid; page-break-inside: avoid; }
            body { background: #fff !important; }
        """)

        # 生成 PDF
        page.pdf(
            path=str(output_path),
            format="A4",
            print_background=True,
            margin={"top": "12mm", "bottom": "12mm", "left": "10mm", "right": "10mm"},
        )
        browser.close()

    return str(output_path), None


# ============================================================
# PPT 导出 — 全板块可编辑文本
# ============================================================

def _esc(s):
    return (s or '').replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


def _sent_type(a):
    s = a.get('sentiment', {})
    if not isinstance(s, dict):
        return 'neutral'
    if s.get('type'):
        return s['type']
    label = str(s.get('sentiment', s.get('label', '')))
    if '利好' in label or '偏多' in label:
        return 'bull'
    if '利空' in label or '偏空' in label:
        return 'bear'
    return 'neutral'


def export_ppt(data_file, desktop):
    """生成可编辑 PPT — 所有元素都是文本框，包含全部板块"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return None, "python-pptx 未安装，请运行: pip install python-pptx"

    data = json.loads(Path(data_file).read_text(encoding='utf-8'))
    articles = data.get('articles', [])
    meta = data.get('meta', {})

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色
    BG = RGBColor(10, 14, 26)
    CARD = RGBColor(17, 24, 39)
    ACCENT = RGBColor(59, 130, 246)
    WHITE = RGBColor(248, 250, 252)
    GRAY = RGBColor(148, 163, 184)
    GREEN = RGBColor(16, 185, 129)
    RED = RGBColor(239, 68, 68)
    YELLOW = RGBColor(245, 158, 11)
    PURPLE = RGBColor(139, 92, 246)
    CYAN = RGBColor(6, 182, 212)
    DARK = RGBColor(71, 85, 105)

    def set_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

    def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
                 bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.font.name = 'Microsoft YaHei'
        return box

    def add_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_table(slide, left, top, headers, rows, col_widths):
        """添加可编辑表格"""
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                            Inches(sum(col_widths)), Inches(0.35 * n_rows))
        tbl = tbl_shape.table
        # 设置列宽
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
        # 表头
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = ACCENT
                p.font.bold = True
                p.font.name = 'Microsoft YaHei'
        # 数据行
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD if ri % 2 == 0 else RGBColor(15, 23, 42)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = WHITE if ci < 2 else GRAY
                    p.font.name = 'Microsoft YaHei'
        return tbl_shape

    # ── 计算统计 ──
    total = len(articles) or 1
    bull = len([a for a in articles if _sent_type(a) == 'bull'])
    bear = len([a for a in articles if _sent_type(a) == 'bear'])
    neu = total - bull - bear
    china = len([a for a in articles if '中国动态' in str(a.get('tags', []))])
    sources = len(set(a.get('source', '') for a in articles))

    # ════════════════════════════════════════════════════════
    # Slide 1: 封面
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 1, 1.5, 11, 1.2, 'AI 行业周报', 52, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, 1, 3.0, 11, 0.6, '多源聚合 · 真实数据 · 情感分析 · 趋势洞察 · 数据看板', 20, GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, 1, 3.8, 11, 0.5, f"{meta.get('week_start', '')} ~ {meta.get('week_end', '')}", 16, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, 1, 4.6, 11, 0.5, f'共 {total} 条新闻 · {sources} 个数据源', 14, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, 1, 5.2, 11, 0.4, f'生成于 {datetime.now().strftime("%Y-%m-%d %H:%M")}', 12, DARK, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # Slide 2: KPI 概览
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '📊 本周 KPI 概览', 30, WHITE, True)

    kpis = [(str(total), '本周新闻', ACCENT), (str(sources), '数据来源', PURPLE),
            (str(bull), '利好信号', GREEN), (str(bear), '利空信号', RED), (str(china), '中国动态', CYAN)]
    cw, gap = 2.2, 0.3
    sx = (13.333 - len(kpis) * cw - (len(kpis) - 1) * gap) / 2
    for i, (num, lbl, clr) in enumerate(kpis):
        x = sx + i * (cw + gap)
        add_rect(slide, x, 1.3, cw, 1.8, CARD)
        add_text(slide, x, 1.5, cw, 0.9, num, 40, clr, True, PP_ALIGN.CENTER)
        add_text(slide, x, 2.5, cw, 0.4, lbl, 13, GRAY, False, PP_ALIGN.CENTER)

    # 情感分布条
    add_text(slide, 0.8, 3.8, 11, 0.4, '情感分布', 18, WHITE, True)
    bw = bull / total * 10
    nw = neu / total * 10
    bearw = bear / total * 10
    if bw > 0:
        add_rect(slide, 0.8, 4.4, max(bw, 0.1), 0.5, GREEN)
    if nw > 0:
        add_rect(slide, 0.8 + bw, 4.4, max(nw, 0.1), 0.5, YELLOW)
    if bearw > 0:
        add_rect(slide, 0.8 + bw + nw, 4.4, max(bearw, 0.1), 0.5, RED)
    add_text(slide, 0.8, 5.1, 11, 0.4,
             f'利好 {bull} ({round(bull/total*100)}%)   中性 {neu} ({round(neu/total*100)}%)   利空 {bear} ({round(bear/total*100)}%)',
             12, GRAY)

    # ════════════════════════════════════════════════════════
    # Slide 3+: 新闻（每页6条）
    # ════════════════════════════════════════════════════════
    per_page = 6
    for pi in range(0, len(articles), per_page):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        page_arts = articles[pi:pi + per_page]
        add_text(slide, 0.5, 0.2, 12, 0.5, f'📰 AI 新闻 ({pi+1}-{min(pi+per_page, len(articles))} / {len(articles)})', 24, WHITE, True)
        y = 0.9
        for a in page_arts:
            st = _sent_type(a)
            clr = GREEN if st == 'bull' else RED if st == 'bear' else YELLOW
            sent = a.get('sentiment', {})
            sent_label = str(sent.get('sentiment', sent.get('label', '中性')) if isinstance(sent, dict) else '中性')
            sent_label = sent_label.replace('🟢 ', '').replace('🔴 ', '').replace('🟡 ', '')

            add_rect(slide, 0.6, y, 0.06, 0.9, clr)
            add_text(slide, 0.85, y, 9.5, 0.35, a.get('title', '')[:80], 14, WHITE, True)
            add_text(slide, 0.85, y + 0.32, 9.5, 0.28, (a.get('summary', '') or '')[:120], 10, GRAY)
            tags = ', '.join(a.get('tags', [])[:3])
            meta_text = f"{a.get('source', '')} · {a.get('published', '')[:16]}"
            if tags:
                meta_text += f'  [{tags}]'
            add_text(slide, 0.85, y + 0.62, 9.5, 0.25, meta_text, 9, DARK)
            add_text(slide, 10.5, y + 0.1, 2.5, 0.3, sent_label, 11, clr, True, PP_ALIGN.RIGHT)
            y += 1.05

    # ════════════════════════════════════════════════════════
    # 情感分析页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '🧠 情感分析', 28, WHITE, True)
    add_text(slide, 0.5, 0.9, 12, 0.4, '基于多维度关键词权重 + 上下文语义规则自动标注', 13, GRAY)

    # 情感卡片
    sent_items = [
        (bull, round(bull/total*100), '利好/偏多', GREEN),
        (neu, round(neu/total*100), '中性', YELLOW),
        (bear, round(bear/total*100), '利空/偏空', RED),
    ]
    sx = 0.8
    for cnt, pct, label, clr in sent_items:
        add_rect(slide, sx, 1.6, 3.5, 1.5, CARD)
        add_text(slide, sx + 0.2, 1.8, 3, 0.7, str(cnt), 36, clr, True, PP_ALIGN.CENTER)
        add_text(slide, sx + 0.2, 2.6, 3, 0.4, f'{label} ({pct}%)', 12, GRAY, False, PP_ALIGN.CENTER)
        sx += 4

    # 情感详情（每类前3条）
    y = 3.5
    for st_type, st_label, st_clr in [('bull', '利好', GREEN), ('bear', '利空', RED)]:
        arts = [a for a in articles if _sent_type(a) == st_type][:3]
        if arts:
            add_text(slide, 0.8, y, 2, 0.3, st_label, 14, st_clr, True)
            for a in arts:
                add_text(slide, 0.8, y + 0.3, 11, 0.25, f"• {a.get('title', '')[:60]}", 10, WHITE)
                y += 0.28
            y += 0.3

    # ════════════════════════════════════════════════════════
    # 周度对比页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '📊 周度对比', 28, WHITE, True)
    add_text(slide, 0.5, 0.9, 12, 0.4, '本周 vs 上周：情感变化、话题热度迁移', 13, GRAY)

    # 尝试加载上周数据
    history_dir = Path(data_file).parent / "history"
    last_week_arts = []
    if history_dir.exists():
        files = sorted(history_dir.glob("*.json"), reverse=True)
        current_week = f"{datetime.now().strftime('%Y')}-W{datetime.now().strftime('%V')}"
        for f in files:
            if f.stem != current_week:
                try:
                    last_week_arts = json.loads(f.read_text('utf-8')).get('articles', [])
                    break
                except:
                    pass

    if last_week_arts:
        lw_total = len(last_week_arts) or 1
        lw_bull = len([a for a in last_week_arts if _sent_type(a) == 'bull'])
        lw_bear = len([a for a in last_week_arts if _sent_type(a) == 'bear'])

        # 情感对比表
        headers = ['指标', '本周', '上周', '变化']
        rows = [
            ['新闻总数', str(total), str(lw_total), f"{'↑' if total >= lw_total else '↓'}{abs(total - lw_total)}"],
            ['利好信号', str(bull), str(lw_bull), f"{'↑' if bull >= lw_bull else '↓'}{abs(bull - lw_bull)}"],
            ['利空信号', str(bear), str(lw_bear), f"{'↑' if bear >= lw_bear else '↓'}{abs(bear - lw_bear)}"],
        ]
        add_table(slide, 0.8, 1.6, headers, rows, [2.5, 2, 2, 2])

        # 话题对比
        cur_tags = {}
        for a in articles:
            for t in a.get('tags', []):
                cur_tags[t] = cur_tags.get(t, 0) + 1
        lw_tags = {}
        for a in last_week_arts:
            for t in a.get('tags', []):
                lw_tags[t] = lw_tags.get(t, 0) + 1
        top_tags = sorted(set(list(cur_tags.keys()) + list(lw_tags.keys())),
                          key=lambda t: cur_tags.get(t, 0), reverse=True)[:6]

        y = 4.0
        add_text(slide, 0.8, y, 3, 0.3, '话题热度变化', 14, ACCENT, True)
        y += 0.4
        for tag in top_tags:
            c = cur_tags.get(tag, 0)
            l = lw_tags.get(tag, 0)
            add_text(slide, 0.8, y, 2, 0.25, tag, 10, WHITE)
            add_text(slide, 3, y, 1.5, 0.25, f'本周: {c}', 10, ACCENT)
            add_text(slide, 4.5, y, 1.5, 0.25, f'上周: {l}', 10, PURPLE)
            diff = c - l
            diff_text = f"{'↑' if diff >= 0 else '↓'}{abs(diff)}"
            add_text(slide, 6, y, 1, 0.25, diff_text, 10, GREEN if diff >= 0 else RED)
            y += 0.3
    else:
        add_text(slide, 0.8, 2, 11, 1, '暂无上周数据\n下次更新时系统会自动保存当前数据，届时可进行周度对比', 16, GRAY, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # 趋势信号页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '📈 趋势信号', 28, WHITE, True)
    trends = [
        (GREEN, '强烈看多', 'AI Agent 从概念走向企业规模化部署', 'Gartner预测40%企业应用将嵌入Agent'),
        (GREEN, '看多', '开源模型性价比碾压闭源', 'DeepSeek V4成本仅为GPT-5的1/30'),
        (GREEN, '看多', '中国AI调用量全球领跑', 'GLM-5完全基于华为昇腾训练'),
        (YELLOW, '中性', '企业端爆发 vs 个人端涨价', '行业从流量争夺转向价值收割'),
        (RED, '风险', 'AI电力危机: 物理基础设施成最大瓶颈', 'Gartner预计2027年40%受电力限制'),
        (RED, '风险', '资本市场情绪转变: 机构资金流出', '机构资金二季度起持续流出'),
    ]
    ty = 1.2
    for clr, label, title, desc in trends:
        add_rect(slide, 0.7, ty, 0.06, 0.8, clr)
        add_text(slide, 0.95, ty, 2, 0.3, label, 12, clr, True)
        add_text(slide, 0.95, ty + 0.3, 11, 0.3, title, 13, WHITE, True)
        add_text(slide, 0.95, ty + 0.55, 11, 0.25, desc, 10, GRAY)
        ty += 1.0

    # ════════════════════════════════════════════════════════
    # 传导链页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '🔗 市场传导链', 28, WHITE, True)
    chains = [
        ('宏观政策', [('AI 监管加强', '政策层', RED), ('合规成本上升', '产业层', PURPLE), ('安全厂商受益', '资本层', GREEN)]),
        ('资本流动', [('AI融资创纪录', '政策层', RED), ('头部公司IPO潮', '产业层', PURPLE), ('芯片供应链受益', '资本层', GREEN)]),
        ('芯片供应', [('芯片出口管制', '政策层', RED), ('国产替代加速', '产业层', PURPLE), ('国产芯片厂受益', '资本层', GREEN)]),
    ]
    cy = 1.2
    for name, nodes in chains:
        add_text(slide, 0.8, cy, 3, 0.3, name, 14, ACCENT, True)
        nx = 0.8
        for title, layer, clr in nodes:
            add_rect(slide, nx, cy + 0.4, 2.8, 0.7, CARD)
            add_text(slide, nx + 0.1, cy + 0.45, 2.6, 0.3, title, 11, WHITE, True, PP_ALIGN.CENTER)
            add_text(slide, nx + 0.1, cy + 0.75, 2.6, 0.25, layer, 9, clr, False, PP_ALIGN.CENTER)
            if nx > 0.8:
                add_text(slide, nx - 0.4, cy + 0.6, 0.4, 0.3, '→', 18, ACCENT, False, PP_ALIGN.CENTER)
            nx += 3.2
        cy += 1.5

    # ════════════════════════════════════════════════════════
    # 价格对比页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '💰 主流模型价格对比', 28, WHITE, True)
    headers = ['厂商', '模型', '输入 $/M tok', '输出 $/M tok', '上下文', '备注']
    prices = [
        ['xAI', 'Grok 4.1', '$0.20', '$0.50', '-', '最便宜前沿模型'],
        ['DeepSeek', 'V3.2', '$0.27', '$1.10', '1M+', '性价比之王'],
        ['OpenAI', 'o4-mini', '$1.10', '$4.40', '-', ''],
        ['Google', 'Gemini 3.1 Pro', '~$1.25', '~$10.00', '2M', '推理翻倍'],
        ['OpenAI', 'GPT-5', '$1.25', '$10.00', '400K', ''],
        ['Anthropic', 'Sonnet 4.6', '$3.00', '$15.00', '1M', '中端逆袭'],
        ['Anthropic', 'Opus 4.6', '$15.00', '$75.00', '200K', '最贵旗舰'],
    ]
    add_table(slide, 0.8, 1.2, headers, prices, [1.8, 2.2, 1.8, 1.8, 1.5, 2.5])

    # ════════════════════════════════════════════════════════
    # 数据看板页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '📊 数据看板', 28, WHITE, True)

    # 来源分布
    src_count = {}
    for a in articles:
        s = a.get('source', '')
        src_count[s] = src_count.get(s, 0) + 1
    src_sorted = sorted(src_count.items(), key=lambda x: -x[1])[:6]
    src_max = max(v for _, v in src_sorted) if src_sorted else 1

    add_text(slide, 0.8, 1.1, 3, 0.3, '来源分布', 14, ACCENT, True)
    sy = 1.5
    for name, cnt in src_sorted:
        add_text(slide, 0.8, sy, 2, 0.25, name, 10, WHITE)
        bar_w = cnt / src_max * 4
        add_rect(slide, 3, sy + 0.05, bar_w, 0.18, ACCENT)
        add_text(slide, 3 + bar_w + 0.1, sy, 1, 0.25, str(cnt), 10, GRAY)
        sy += 0.32

    # 标签分布
    tag_count = {}
    for a in articles:
        for t in a.get('tags', []):
            tag_count[t] = tag_count.get(t, 0) + 1
    tag_sorted = sorted(tag_count.items(), key=lambda x: -x[1])[:6]
    tag_max = max(v for _, v in tag_sorted) if tag_sorted else 1

    add_text(slide, 6.5, 1.1, 3, 0.3, '话题标签分布', 14, PURPLE, True)
    ty = 1.5
    for name, cnt in tag_sorted:
        add_text(slide, 6.5, ty, 2, 0.25, name, 10, WHITE)
        bar_w = cnt / tag_max * 4
        add_rect(slide, 8.7, ty + 0.05, bar_w, 0.18, PURPLE)
        add_text(slide, 8.7 + bar_w + 0.1, ty, 1, 0.25, str(cnt), 10, GRAY)
        ty += 0.32

    # 高频关键词
    add_text(slide, 0.8, 4.0, 3, 0.3, '高频关键词', 14, CYAN, True)
    keywords = ['AI', 'Agent', '大模型', '芯片', '融资', '开源', 'DeepSeek', 'OpenAI', 'Google', 'Claude']
    kw_text = '  ·  '.join(keywords)
    add_text(slide, 0.8, 4.4, 11, 0.5, kw_text, 14, GRAY)

    # 中国 vs 全球
    add_text(slide, 0.8, 5.2, 3, 0.3, '中国 vs 全球', 14, ACCENT, True)
    glb = total - china
    add_text(slide, 0.8, 5.6, 5, 0.3, f'中国: {china} ({round(china/total*100)}%)', 12, ACCENT)
    add_text(slide, 6, 5.6, 5, 0.3, f'全球: {glb} ({round(glb/total*100)}%)', 12, PURPLE)

    # ════════════════════════════════════════════════════════
    # 时间线页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 0.3, 12, 0.6, '📅 时间线', 28, WHITE, True)
    tl = sorted(articles, key=lambda x: x.get('published', ''), reverse=True)[:10]
    ty = 1.0
    for a in tl:
        add_rect(slide, 1.0, ty + 0.1, 0.12, 0.12, ACCENT)
        add_text(slide, 1.3, ty - 0.05, 3, 0.25, a.get('published', '')[:16], 10, ACCENT, True)
        add_text(slide, 1.3, ty + 0.2, 10, 0.3, a.get('title', ''), 12, WHITE, True)
        add_text(slide, 1.3, ty + 0.5, 10, 0.25, (a.get('summary', '') or '')[:100], 9, GRAY)
        ty += 0.85

    # ════════════════════════════════════════════════════════
    # 尾页
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 1, 2.2, 11, 1, 'AI 行业周报', 48, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, 1, 3.5, 11, 0.5, '数据来源:', 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 1, 4.0, 11, 0.5, 'TechCrunch · The Verge · Wired · VentureBeat · AI News · MIT科技评论 · 纽约时报 · BBC', 12, GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, 1, 4.5, 11, 0.5, 'IT之家 · 机器之心 · 36氪 · 虎嗅 · arXiv · OpenAI · Anthropic · Google · Meta · Microsoft', 12, GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, 1, 5.2, 11, 0.5, '内容由 AI 自动聚合，仅供参考', 12, DARK, False, PP_ALIGN.CENTER)

    # 保存
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = Path(desktop) / f"AI行业周报_{timestamp}.pptx"
    prs.save(str(output_path))
    return str(output_path), None
